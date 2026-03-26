import os
from pathlib import Path
from uuid import uuid4


TMP_DIR = Path("tmp")
RUN_ROOT = TMP_DIR / "real_qwen_runs" / uuid4().hex[:8]
os.environ["APP_DATABASE_PATH"] = str(RUN_ROOT / "app.db")
os.environ["APP_UPLOAD_ROOT"] = str(RUN_ROOT / "uploads")
os.environ["APP_OUTPUT_ROOT"] = str(RUN_ROOT / "outputs")

import fitz
from fastapi.testclient import TestClient
from pptx import Presentation

from backend.app.main import app


def create_sample_pptx() -> Path:
    RUN_ROOT.mkdir(parents=True, exist_ok=True)
    pptx_path = RUN_ROOT / "real_probe_courseware.pptx"

    presentation = Presentation()

    slide1 = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide1.shapes.title.text = "监督学习概述"
    slide1.placeholders[1].text = (
        "监督学习使用带标签数据训练模型。\n"
        "目标是学习输入到输出的映射关系。\n"
        "常见任务包括分类和回归。"
    )

    slide2 = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide2.shapes.title.text = "分类与回归"
    slide2.placeholders[1].text = (
        "分类任务输出离散类别，例如垃圾邮件识别。\n"
        "回归任务输出连续数值，例如房价预测。\n"
        "标签为模型提供正确答案，帮助计算误差并优化参数。"
    )

    presentation.save(pptx_path)
    return pptx_path


def main() -> None:
    pptx_path = create_sample_pptx()

    with TestClient(app) as client:
        with pptx_path.open("rb") as handle:
            upload_response = client.post(
                "/api/documents/upload",
                files={
                    "file": (
                        pptx_path.name,
                        handle,
                        "application/vnd.openxmlformats-officedocument.presentationml.presentation",
                    )
                },
            )
        upload_payload = upload_response.json()
        document = upload_payload["document"]
        document_id = document["id"]

        print("upload_status", upload_response.status_code)
        print("document_title", document["title"])
        print("sections", len(document["sections"]))
        print("overview_len", len(document["learning_board"]["overview"]))

        first_question = "请详细讲解监督学习，并说明分类和回归的区别。"
        first_chat = client.post(
            f"/api/documents/{document_id}/chat",
            json={"question": first_question, "scope": "all"},
        )
        first_payload = first_chat.json()
        print("chat1_status", first_chat.status_code)
        print("chat1_len", len(first_payload["answer"]))
        print("chat1_answer", first_payload["answer"])

        second_question = "那它为什么需要标签？"
        second_chat = client.post(
            f"/api/documents/{document_id}/chat",
            json={"question": second_question, "scope": "all"},
        )
        second_payload = second_chat.json()
        print("chat2_status", second_chat.status_code)
        print("chat2_len", len(second_payload["answer"]))
        print("chat2_answer", second_payload["answer"])


if __name__ == "__main__":
    main()
