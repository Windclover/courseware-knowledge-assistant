import os
from pathlib import Path
from uuid import uuid4


TMP_DIR = Path("tmp")
SMOKE_ROOT = TMP_DIR / "smoke_runs" / uuid4().hex[:8]
os.environ["APP_DATABASE_PATH"] = str(SMOKE_ROOT / "app.db")
os.environ["APP_UPLOAD_ROOT"] = str(SMOKE_ROOT / "uploads")
os.environ["APP_OUTPUT_ROOT"] = str(SMOKE_ROOT / "outputs")

import fitz
from fastapi.testclient import TestClient
from pptx import Presentation

from backend.app.domain import ChatAnswerResult, GeneratedSection
import backend.app.main as main_module
import backend.app.services.pipeline as pipeline_module
from backend.app.schemas import (
    AssessmentQuestion,
    AssessmentQuestionOption,
    AssessmentSuite,
    FormulaNote,
    LearningBoard,
    LearningConceptItem,
    LearningPracticeItem,
    LearningReviewStep,
    LearningSummaryItem,
    WorkedExample,
)


class FakeQwen:
    def generate_section(self, document_title, draft):
        return GeneratedSection(
            section_index=draft.section_index,
            title=draft.title,
            detailed_explanation=f"{draft.title} 的详细讲解：用于复习与答辩展示。",
            key_points=[
                "监督学习依赖带标签数据",
                "分类与回归是典型任务",
            ],
            formula_notes=[
                FormulaNote(
                    title="监督学习映射关系",
                    latex="y = f(x)",
                    raw_text="",
                    explanation="监督学习学习输入到输出的映射关系。",
                    source_refs=draft.source_refs,
                )
            ],
            worked_examples=[
                WorkedExample(
                    title="分类与回归例题",
                    problem="判断垃圾邮件识别和房价预测分别属于哪类任务。",
                    steps=["识别输出类型。", "离散类别对应分类，连续数值对应回归。"],
                    final_answer="垃圾邮件识别是分类，房价预测是回归。",
                    source_refs=draft.source_refs,
                )
            ],
            source_refs=draft.source_refs,
        )

    def generate_learning_board(self, *, document_title, sections):
        return LearningBoard(
            overview=f"《{document_title}》的复习重点是先理解监督学习，再区分分类与回归。",
            summary=[
                LearningSummaryItem(
                    id="summary-1",
                    title=sections[0].title,
                    summary="监督学习依赖带标签数据。",
                    source_refs=sections[0].source_refs,
                )
            ],
            concepts=[
                LearningConceptItem(
                    id="concept-1",
                    term="监督学习",
                    explanation="监督学习通过带标签样本建立输入到输出的映射。",
                    section_title=sections[0].title,
                    source_refs=sections[0].source_refs,
                )
            ],
            practice=[
                LearningPracticeItem(
                    id="practice-1",
                    prompt="判断一个案例属于分类还是回归。",
                    section_title=sections[0].title,
                    source_refs=sections[0].source_refs,
                )
            ],
            review_path=[
                LearningReviewStep(
                    id="review-1",
                    title="先看摘要再看概念",
                    detail="先读摘要，再理解监督学习与分类/回归的关系。",
                    source_refs=sections[0].source_refs,
                )
            ],
        )

    def generate_assessment(self, *, document_title, sections, learning_board):
        return AssessmentSuite(
            title=f"{document_title} 测试环境",
            intro="完成学习任务后开始答题。",
            questions=[
                AssessmentQuestion(
                    id="q1",
                    type="choice",
                    prompt="监督学习依赖什么数据？",
                    options=[
                        AssessmentQuestionOption(id="A", text="带标签数据"),
                        AssessmentQuestionOption(id="B", text="随机噪声"),
                    ],
                    answer="A",
                    acceptable_answers=[],
                    explanation="监督学习通过带标签样本学习输入到输出的映射。",
                ),
                AssessmentQuestion(
                    id="q2",
                    type="calculation",
                    prompt="分类与____是两类典型监督学习任务。",
                    options=[],
                    answer="回归",
                    display_answer="回归",
                    acceptable_answers=["回归"],
                    explanation="课件中明确指出分类与回归是典型监督学习任务。",
                    solution_steps=["先识别监督学习的常见任务。", "课件明确给出分类与回归。"],
                ),
            ],
        )

    def answer_question(self, **kwargs):
        return ChatAnswerResult(
            answer="这是基于课件内容生成的回答。",
            source_refs=["第 1 张幻灯片"],
            supplemental=True,
            supplemental_notes="这里补充了一点通用常识。",
        )


def create_sample_files() -> tuple[Path, Path]:
    SMOKE_ROOT.mkdir(parents=True, exist_ok=True)

    pptx_path = SMOKE_ROOT / "sample_courseware.pptx"
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = "监督学习"
    slide.placeholders[1].text = "监督学习依赖带标签数据\n常见任务包括分类和回归"
    presentation.save(pptx_path)

    pdf_path = SMOKE_ROOT / "sample_courseware.pdf"
    document = fitz.open()
    page = document.new_page()
    page.insert_text(
        (72, 72),
        "监督学习导论\n监督学习使用带标签数据。\n分类与回归是典型任务。",
    )
    document.save(pdf_path)
    document.close()

    return pptx_path, pdf_path


def main() -> None:
    pipeline_module.QwenClient = FakeQwen
    main_module.QwenClient = FakeQwen
    pptx_path, pdf_path = create_sample_files()

    with TestClient(main_module.app) as client:
        last_document_id = ""
        for file_path, content_type in [
            (
                pptx_path,
                "application/vnd.openxmlformats-officedocument.presentationml.presentation",
            ),
            (pdf_path, "application/pdf"),
        ]:
            with file_path.open("rb") as handle:
                response = client.post(
                    "/api/documents/upload",
                    files={"file": (file_path.name, handle, content_type)},
                )
            payload = response.json()
            if response.status_code != 200:
                print("upload_failed", file_path.name, response.status_code, payload)
                return
            last_document_id = payload["document"]["id"]
            print(
                file_path.name,
                response.status_code,
                payload["document"]["sections"][0]["title"],
                len(payload["document"]["sources"]),
                len(payload["document"]["learning_board"]["concepts"]),
                len(payload["document"]["learning_board"]["practice"]),
            )

        chat_response = client.post(
            f"/api/documents/{last_document_id}/chat",
            json={"question": "这一章讲了什么？", "scope": "all"},
        )
        chat_payload = chat_response.json()
        print("chat", chat_response.status_code, chat_payload["source_refs"])

        assessment_response = client.post(
            f"/api/documents/{last_document_id}/assessment"
        )
        assessment_payload = assessment_response.json()
        print(
            "assessment",
            assessment_response.status_code,
            len(assessment_payload["questions"]),
        )


if __name__ == "__main__":
    main()
