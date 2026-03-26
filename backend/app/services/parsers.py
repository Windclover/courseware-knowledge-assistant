from __future__ import annotations

from collections import Counter
from pathlib import Path
import re

import fitz
from pptx import Presentation

from ..domain import ParsedFragment


class UnsupportedDocumentError(ValueError):
    pass


def parse_document(file_path: Path) -> tuple[str, str, list[ParsedFragment]]:
    suffix = file_path.suffix.lower()
    if suffix == ".pptx":
        return _parse_pptx(file_path)
    if suffix == ".pdf":
        return _parse_pdf(file_path)
    raise UnsupportedDocumentError("仅支持 PPTX 或文本型 PDF。")


def _parse_pptx(file_path: Path) -> tuple[str, str, list[ParsedFragment]]:
    presentation = Presentation(file_path)
    document_title = presentation.core_properties.title or file_path.stem
    fragments: list[ParsedFragment] = []

    for index, slide in enumerate(presentation.slides, start=1):
        title = None
        if slide.shapes.title and slide.shapes.title.text:
            title = _clean_text(slide.shapes.title.text)
        texts: list[str] = []
        for shape in slide.shapes:
            if not getattr(shape, "has_text_frame", False):
                continue
            text = _clean_text(shape.text)
            if not text:
                continue
            if title and text == title:
                continue
            texts.append(text)
        slide_text = "\n".join(texts).strip()
        if not slide_text and title:
            slide_text = title
        if not slide_text:
            continue
        fragments.append(
            ParsedFragment(
                fragment_index=index,
                source_label=f"第 {index} 张幻灯片",
                source_type="slide",
                title=title,
                text=slide_text,
            )
        )

    if not fragments:
        raise UnsupportedDocumentError("PPTX 中未提取到有效文本。")
    return document_title, "pptx", fragments


def _parse_pdf(file_path: Path) -> tuple[str, str, list[ParsedFragment]]:
    document = fitz.open(file_path)
    metadata_title = (document.metadata or {}).get("title") if document.metadata else None
    document_title = metadata_title or file_path.stem
    fragments: list[ParsedFragment] = []
    page_lines = [_extract_pdf_lines(page.get_text("text")) for page in document]
    noise_lines = _detect_pdf_noise_lines(page_lines, document.page_count)

    for index, lines in enumerate(page_lines, start=1):
        filtered_lines = _filter_pdf_lines(lines, noise_lines)
        if _should_skip_pdf_page(filtered_lines):
            continue
        text = "\n".join(filtered_lines).strip()
        if not text:
            continue
        title = _detect_pdf_page_title(filtered_lines)
        fragments.append(
            ParsedFragment(
                fragment_index=index,
                source_label=f"第 {index} 页",
                source_type="page",
                title=title,
                text=text,
            )
        )

    document.close()
    if not fragments:
        raise UnsupportedDocumentError(
            "未检测到可提取文本，可能是扫描版 PDF。首版暂不支持 OCR。"
        )
    return document_title, "pdf", fragments


def _extract_pdf_lines(text: str) -> list[str]:
    lines: list[str] = []
    for raw_line in text.splitlines():
        line = _clean_pdf_line(raw_line)
        if line:
            lines.append(line)
    return lines


def _clean_pdf_line(text: str) -> str:
    normalized = text.replace("\u3000", " ").replace("\x0b", " ")
    normalized = re.sub(r"[ \t]+", " ", normalized)
    return normalized.strip()


def _detect_pdf_noise_lines(
    page_lines: list[list[str]],
    page_count: int,
) -> set[str]:
    counts: Counter[str] = Counter()
    threshold = max(6, page_count // 3)
    for lines in page_lines:
        seen: set[str] = set()
        for line in lines:
            normalized = _normalize_pdf_line(line)
            if normalized:
                seen.add(normalized)
        counts.update(seen)

    return {
        line
        for line, count in counts.items()
        if _looks_like_footer_or_page_number(line)
        or (count >= threshold and not _looks_like_catalog_item(line))
    }


def _filter_pdf_lines(lines: list[str], noise_lines: set[str]) -> list[str]:
    filtered: list[str] = []
    for line in lines:
        normalized = _normalize_pdf_line(line)
        if not normalized:
            continue
        if normalized in noise_lines:
            continue
        if _is_garbled_catalog_line(line):
            continue
        filtered.append(line)
    return filtered


def _detect_pdf_page_title(lines: list[str]) -> str | None:
    candidates: list[str] = []
    for line in lines[:12]:
        if _is_title_candidate(line):
            candidates.append(line)
            continue
        if candidates and (_is_formula_like_line(line) or len(line) > 18 or line.startswith("•")):
            break
    if candidates:
        return candidates[-1]
    for line in lines[:8]:
        if _is_title_candidate(line):
            return line
    return None


def _should_skip_pdf_page(lines: list[str]) -> bool:
    if not lines:
        return True
    text = "\n".join(lines)
    if "副教授" in text or "@" in text:
        return True
    catalog_hits = sum(1 for line in lines if _looks_like_catalog_item(line))
    if lines[0] == "目录":
        return True
    if catalog_hits >= 4 and not any(
        _is_formula_like_line(line) or line.startswith("•") or len(line) > 18
        for line in lines[4:]
    ):
        return True
    if len(text) < 20 and not any(_is_formula_like_line(line) for line in lines):
        return True
    return False


def _normalize_pdf_line(text: str) -> str:
    return re.sub(r"\s+", " ", text.strip())


def _looks_like_footer_or_page_number(text: str) -> bool:
    normalized = _normalize_pdf_line(text)
    return bool(
        re.fullmatch(r"\d+\s*/\s*\d+", normalized)
        or normalized in {"最优化方法", "Optimization Methods", "目录"}
    )


def _is_garbled_catalog_line(text: str) -> bool:
    normalized = _normalize_pdf_line(text)
    if not normalized:
        return True
    if len(normalized) <= 2 and re.fullmatch(r"[\u4e00-\u9fff]+", normalized):
        return True
    return False


def _looks_like_cover_line(text: str) -> bool:
    return any(token in text for token in ("副教授", "@", "中国农业大学"))


def _looks_like_catalog_item(text: str) -> bool:
    candidates = (
        "罚函数法的含义",
        "外点罚函数方法",
        "等式约束",
        "不等式约束",
        "一般最优化问题",
        "内点罚函数方法",
    )
    return any(token in text for token in candidates)


def _is_formula_like_line(text: str) -> bool:
    math_tokens = ("=", "min", "argmin", "s.t", "∈", "∇", "σ", "λ", "μ", "||", "⩽")
    if any(token in text for token in math_tokens):
        return True
    if re.search(r"[A-Za-z]\w*\(.*\)", text):
        return True
    if re.search(r"x\d|c\d|g\d|h\d", text):
        return True
    return False


def _is_title_candidate(text: str) -> bool:
    if _is_formula_like_line(text):
        return False
    if _looks_like_footer_or_page_number(text):
        return False
    if _looks_like_cover_line(text):
        return False
    if text.startswith("•") or text.endswith("："):
        return False
    if any(token in text for token in ("，", "。", "；")) and len(text) > 18:
        return False
    return bool(len(text) <= 28 and re.search(r"[\u4e00-\u9fffA-Za-z]", text))


def _clean_text(text: str) -> str:
    normalized = text.replace("\x0b", "\n").replace("\u3000", " ")
    normalized = re.sub(r"\n{3,}", "\n\n", normalized)
    normalized = re.sub(r"[ \t]+", " ", normalized)
    return normalized.strip()
